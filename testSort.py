#!/usr/bin/env python
# _*_ coding:utf-8 _*_
import numpy as np
from pptx import Presentation
from pptx_tools import utils
import os

global cwd
cwd = os.getcwd()


# pptdir = cwd + r'\ecd.pptx'
# pptx = Presentation(pptdir)
# # 遍历ppt文件的所有幻灯片页
# for slide in pptx.slides:
#     # 遍历幻灯片页的所有形状
#     print(1)
#     for shape in slide.shapes:
#     # 判断形状是否含有文本框，如果含有则顺序运行代码
#         if shape.has_text_frame:
#             # 获取文本框
#             text_frame = shape.text_frame
#             # 遍历文本框中的所有段落
#             for paragraph in text_frame.paragraphs:
#                 # 将文本框中的段落文字写入word中
#                 print(paragraph.text)


prs = Presentation()
title_slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(title_slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]


title.text = 'Hello world!'
subtitle.text = 'python-pptx was here.'

prs.save(cwd + r'\test.pptx')

pptfile = cwd + r'\test.pptx'

# print(pptfile)
png_folder = cwd +r'\temp'
#png_folder = cwd +r'\1.png'
print(png_folder)
utils.save_pptx_as_png(png_folder, pptfile, overwrite_folder=True)

########################################################################################################
# a = np.array([0,1,3,2,6,4,5])
# b = np.array([0,1,2,3,4,5,6])
# index = np.lexsort((b, a))
#
# sorted_index = np.argsort(a)
#
# print(index)
# print(sorted_index)
#
# c  = a[index]
# print(c)
# #name sort by second word
# [i for i,v in sorted(enumerate(['Vincent', 'Alex', 'Bill', 'Matthew']), key=lambda x:x[1])]
#
# x =[1,2,3]
# y =[-4,5,6]
# z = list(map(abs,y))
# print(z)