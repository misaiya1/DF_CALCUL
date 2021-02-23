#!/usr/bin/env python
# _*_ coding:utf-8 _*_
from UI.noname import DF_CALCUL
import numpy as np
from math import pi
import datetime

import os
import sys
import matplotlib.pyplot as plt
import wx
import wx.xrc
import math
from docx import Document
from docx.shared import Inches
from docx.shared import Pt
from docx.shared import RGBColor
from docx.oxml.ns import qn
from docx.enum.text import WD_PARAGRAPH_ALIGNMENT

from docxtpl import DocxTemplate
from pptx import Presentation
from pptx_tools import utils
from pptx.util import Inches, Pt

myChecked = False
NetFreq = 50

''' 函数：返回变量是否被定义过 '''


def isset(v):
    try:
        type(eval(v))
    except:
        return 0
    else:
        return 1


class MyFrame(DF_CALCUL):
    ''' 创建输出路径  '''
    global cwd
    cwd = os.getcwd()
    print(cwd)
    global outputFolderName

    outputFolderName = cwd + r'\fig_save'  # + datetime.datetime.now().strftime('%Y.%m.%d.%H.%M')
    outPathExists = os.path.exists(outputFolderName)
    if not outPathExists:
        os.makedirs(outputFolderName)
        print(outputFolderName + ' 路径创建成功')
    else:
        print(outputFolderName + ' 路径已存在')

    def __init__(self, parent):
        DF_CALCUL.__init__(self, parent)

        wx.Frame.SetBackgroundColour(self, 'Write')  # ???
        font1 = wx.Font(22, wx.MODERN, wx.NORMAL, wx.NORMAL)
        self.m_staticText15.SetFont(font1)
        self.m_button.SetFont(font1)
        self.m_button2.SetFont(font1)
        self.m_staticText16.SetFont(font1)
        self.m_grid2.SetRowLabelValue(0, '转速[RPM]')  # ???
        self.m_grid2.SetRowLabelValue(1, '上网功率[kW]')  # ???
        # ctrl - c

        ########################################################## 3.0 - 150
        self.m_grid2.SetCellValue(0, 0, '1050')  # 第1行第0个
        self.m_grid2.SetCellValue(0, 1, '1050')
        self.m_grid2.SetCellValue(0, 2, '1121')
        self.m_grid2.SetCellValue(0, 3, '1354')
        self.m_grid2.SetCellValue(0, 4, '1494')
        self.m_grid2.SetCellValue(0, 5, '1587')
        self.m_grid2.SetCellValue(0, 6, '1725')
        self.m_grid2.SetCellValue(0, 7, '1725')

        self.m_grid2.SetCellValue(1, 0, '6.76')  # 第1行第0个
        self.m_grid2.SetCellValue(1, 1, '327.508')
        self.m_grid2.SetCellValue(1, 2, '445.227')
        self.m_grid2.SetCellValue(1, 3, '850.095')
        self.m_grid2.SetCellValue(1, 4, '1167.75')
        self.m_grid2.SetCellValue(1, 5, '1410.65')
        self.m_grid2.SetCellValue(1, 6, '1831.74')
        self.m_grid2.SetCellValue(1, 7, '3000')

        self.statusbar = self.CreateStatusBar()
        self.Centre()

    def OnCheck(self, event):
        global myChecked
        w = 2 * pi * NetFreq

        myChecked = self.m_checkBox1.GetValue()
        if myChecked:
            self.m_Lksss.SetLabel(u"定子漏抗[ohm]")
            self.m_Lkrrr.SetLabel(u"转子漏抗[ohm]")
            self.m_staticText22.SetLabel(u"激磁电抗[ohm]")

            self.m_Lm.SetValue(str(float(self.m_Lm.GetValue()) * w))
            self.m_Lks.SetValue(str(float(self.m_Lks.GetValue()) * w))
            self.m_Lkr.SetValue(str(float(self.m_Lkr.GetValue()) * w))



        else:
            self.m_Lksss.SetLabel(u"定子漏感[H]")
            self.m_Lkrrr.SetLabel(u"转子漏感[H]")
            self.m_staticText22.SetLabel(u"电机互感[H]")

            self.m_Lm.SetValue(str(float(self.m_Lm.GetValue()) / w))
            self.m_Lks.SetValue(str(float(self.m_Lks.GetValue()) / w))
            self.m_Lkr.SetValue(str(float(self.m_Lkr.GetValue()) / w))

        print('myChecked   %f' % myChecked)

    def m_buttonOnButtonClick(self, event):
        print("press button1")
        self.DoWork(1)
        self.PlotAll()

    def m_buttonOnButtonClick2(self, event):
        print("press button2")
        self.DoWork(1)
        self.Plot1()
        self.Plot2()
        self.Plot3()

        ###################################################################报错
        # self.Docu()
        # self.statusbar.SetStatusText('Report Generate at , %s ' % wx.Now())
        # dlg = wx.MessageDialog(None, u"报告已生成，请确认", u"完成",
        #                        wx.YES_DEFAULT | wx.ICON_QUESTION)
        # if dlg.ShowModal() == wx.ID_YES:
        #     self.Close(True)
        # dlg.Destroy()
        ###################################################################不报错

        try:
            self.Docu()
            self.statusbar.SetStatusText('Report Generate at , %s ' % wx.Now())
            dlg = wx.MessageDialog(None, u"报告已生成，请确认", u"完成", wx.YES_DEFAULT | wx.ICON_QUESTION)
            if dlg.ShowModal() == wx.ID_YES:
                self.Close(True)
            dlg.Destroy()
        except:
            dlg = wx.MessageDialog(None, u"请先关闭生成物", u"完成", wx.YES_DEFAULT | wx.ICON_QUESTION)
            if dlg.ShowModal() == wx.ID_YES:
                self.Close(True)
            dlg.Destroy()

    def m_slider1OnScrollChanged(self, event):
        pass
        # temp = float(self.m_slider1.GetValue())
        # self.m_textCtrlS.Value = ("%8.3f" % temp)

    def DoWork(self, netVotGain):

        Root2 = 1.414
        Root3 = 1.732
        plt.close('all')
        # slider input part
        # tempMin = float(self.m_textCtrlS1.GetValue())
        # tempMax = float(self.m_textCtrlS2.GetValue())
        # self.m_slider1.SetMin(tempMin)
        # self.m_slider1.SetMax(tempMax)
        # self.m_slider1.Value = temp
        # print(self.m_slider1.Value)

        # input part
        global pp, NetFreq, TongBuZhuanSu, StatorRePower, RotorOpenVoltage, GenRpm, XiaoLv
        global Xm, Xs, Xr, Rs, Rr

        pp = float(self.m_pp.GetValue())  # pp
        A2 = pp
        NetFreq = float(self.m_hz.GetValue())  # 电网频率
        N2 = NetFreq

        TongBuZhuanSu = 60 * NetFreq / pp  # 同步转速
        B2 = TongBuZhuanSu

        # GoNetPower = float(self.m_GoNetPower.GetValue()) * 1E3  # 上网总功率
        # C2 = GoNetPower
        StatorRePower = float(self.m_StatorRePower.GetValue()) * 1E3  # 定子无功
        D2 = StatorRePower
        RotorOpenVoltage = float(self.m_RotorOpen.GetValue())  # 转子开口电压
        E2 = RotorOpenVoltage
        GridVRMS = float(self.m_GridVRMS.GetValue()) * netVotGain  # 电网电压
        F2 = GridVRMS
        G2 = Root2 * GridVRMS / Root3
        print(G2)
        GenRpm = float(self.m_textCtrlS.GetValue())  # 发电机转
        P2 = GenRpm
        XiaoLv = float(self.m_XiaoLv.GetValue())  # 效率
        M2 = XiaoLv

        ###########################################################################
        Rs = float(self.m_Rs.GetValue())  # 定子电阻
        H2 = Rs
        print('Rs = ')
        print(Rs)
        Rr = float(self.m_Rr.GetValue())  # 转子电阻
        I2 = Rr

        # myChecked = self.m_checkBox1.GetValue()
        # if isset('myChecked'):
        tempW = NetFreq * 2 * pi

        if not myChecked:
            Lm = float(self.m_Lm.GetValue())  # 电机互感
            L2 = Lm

            Lks = float(self.m_Lks.GetValue())  # 定子漏感
            J2 = Lks
            Lkr = float(self.m_Lkr.GetValue())  # 转子漏感
            K2 = Lkr
        else:

            Lm = float(self.m_Lm.GetValue()) / tempW  # 电机互感
            L2 = Lm

            Lks = float(self.m_Lks.GetValue()) / tempW  # 定子漏感
            J2 = Lks
            Lkr = float(self.m_Lkr.GetValue()) / tempW  # 转子漏感
            K2 = Lkr
        Xm = Lm * tempW
        Xs = Lks * tempW
        Xr = Lkr * tempW
        print('myChecked   HAVE %f' % myChecked)

        # else:
        #     Lm = float(self.m_Lm.GetValue())  # 电机互感
        #     L2 = Lm
        #     Lks = float(self.m_Lks.GetValue())  # 定子漏感
        #     J2 = Lks
        #     Lkr = float(self.m_Lkr.GetValue())  # 转子漏感
        #     K2 = Lkr
        #     print('myChecked  None ')

        print('======================================')
        print('Lks  %f' % Lks)
        print('Lkr  %f' % Lkr)
        print('Lm  %f' % Lm)

        print('======================================')
        ############################################################################
        global iGenRpm, iGoNetPower, tempMin, tempMax, n
        tempMin = 1e6
        tempMax = 0
        tempLeft = 0
        tempLeft1 = 0
        tempRight = 1e6
        tempRight1 = 1e6
        tempGoNet = 0;
        self.m_grid2.SetRowLabelValue(0, '转速[RPM]')  # ???
        self.m_grid2.SetRowLabelValue(1, '上网功率[kW]')  # ???
        m = 26
        n = 0
        baohan = 0

        iGenRpm = []
        iGoNetPower = []
        for i in range(m):
            if (self.m_grid2.GetCellValue(0, i) == '') or (self.m_grid2.GetCellValue(1, i) == ''):
                pass
                # print('empty %s' % i)
            else:
                # print('full %s' % i)
                temp1 = float(self.m_grid2.GetCellValue(0, i))
                temp2 = float(self.m_grid2.GetCellValue(1, i))

                iGenRpm.append(float(temp1))
                iGoNetPower.append(float(temp2) * 1E3)
                tempMin = min(tempMin, float(self.m_grid2.GetCellValue(0, i)))
                tempMax = max(tempMin, float(self.m_grid2.GetCellValue(0, i)))

                if GenRpm == temp1:  # 包含
                    baohan = 1
                    tempGoNet = temp2

                if temp1 < GenRpm and temp1 >= tempLeft:
                    tempLeft = temp1
                    tempLeft1 = temp2

                if temp1 > GenRpm and temp1 <= tempRight:
                    tempRight = temp1
                    tempRight1 = temp2

                n = n + 1
        print('num %s' % n)

        if GenRpm < tempMin:
            GenRpm = tempMin
            self.m_textCtrlS.SetValue(str(tempMin))

        if GenRpm > tempMax:
            GenRpm = tempMax
            print('tempMax %s' % tempMax)
            self.m_textCtrlS.SetValue(str(tempMax))

        # 队列最后加入 单点值
        iGenRpm.append(GenRpm)
        if baohan == 1:
            iGoNetPower.append(tempGoNet * 1E3)
        else:
            tempGoNet = (GenRpm - tempLeft) / (tempRight - tempLeft) * (tempRight1 - tempLeft1) + tempLeft1
            iGoNetPower.append(tempGoNet * 1E3)

        n = n + 1
        ##########################################################################作图

        iP2 = [0 for i in range(n)]
        iQ2 = [0 for i in range(n)]
        iR2 = [0 for i in range(n)]
        iS2 = [0 for i in range(n)]
        iT2 = [0 for i in range(n)]
        iU2 = [0 for i in range(n)]
        iV2 = [0 for i in range(n)]
        iW2 = [0 for i in range(n)]
        iX2 = [0 for i in range(n)]
        iY2 = [0 for i in range(n)]
        iZ2 = [0 for i in range(n)]
        iAA2 = [0 for i in range(n)]
        iAB2 = [0 for i in range(n)]
        iAC2 = [0 for i in range(n)]
        iAD2 = [0 for i in range(n)]
        iAE2 = [0 for i in range(n)]
        iAF2 = [0 for i in range(n)]

        global iZhuanChaLv
        iZhuanChaLv = [0 for i in range(n)]
        global iStatorAPower
        iStatorAPower = [0 for i in range(n)]
        global iRotorAPower
        iRotorAPower = [0 for i in range(n)]
        global iTotalAPower
        iTotalAPower = [0 for i in range(n)]
        global iRotorVoltRMS
        iRotorVoltRMS = [0 for i in range(n)]
        global iGenId
        iGenId = [0 for i in range(n)]
        global iGenIq
        iGenIq = [0 for i in range(n)]
        global iGenIrmsGS
        iGenIrmsGS = [0 for i in range(n)]
        global iGenIrmsSJ
        iGenIrmsSJ = [0 for i in range(n)]
        global iStatorId
        iStatorId = [0 for i in range(n)]
        global iStatorIq
        iStatorIq = [0 for i in range(n)]
        global iStatorIrms
        iStatorIrms = [0 for i in range(n)]
        global iNetIrms
        iNetIrms = [0 for i in range(n)]
        global iGenVd
        iGenVd = [0 for i in range(n)]
        global iGenVq
        iGenVq = [0 for i in range(n)]
        global iGenV
        iGenV = [0 for i in range(n)]
        global iGenVrms
        iGenVrms = [0 for i in range(n)]
        global iTorque
        iTorque = [0 for i in range(n)]
        global iGridrms
        iGridrms = [0 for i in range(n)]

        for i in range(n):
            if iGenRpm[i] == TongBuZhuanSu:
                iGenRpm[i] = iGenRpm[i] + 1
            iZhuanChaLv[i] = (TongBuZhuanSu - iGenRpm[i]) / TongBuZhuanSu  # 转差率
            iQ2[i] = iZhuanChaLv[i]
            iStatorAPower[i] = iGoNetPower[i] / (1 - iZhuanChaLv[i])  # 定子有功
            iR2[i] = iStatorAPower[i]
            iRotorAPower[i] = iStatorAPower[i] * iZhuanChaLv[i]  # 转子有功
            iS2[i] = iRotorAPower[i]
            iRotorVoltRMS[i] = RotorOpenVoltage * iZhuanChaLv[i]  # 转子电压有效值
            iT2[i] = iRotorVoltRMS[i]
            #########################################################################
            iGenId[i] = iR2[i] / 1.5 / (L2 / (L2 + J2) * G2)
            iU2 = iGenId
            iGenIq[i] = (D2 / 1.5 - G2 ** 2 / (L2 + J2) / 2 / 3.14 / N2) / (L2 * G2 / (L2 + J2))
            iV2[i] = iGenIq[i]

            iGenIrmsGS[i] = (iU2[i] ** 2 + iV2[i] ** 2) ** 0.5 / 1.414
            iAD2[i] = iGenIrmsGS[i]
            iGenIrmsSJ[i] = iAD2[i] * F2 / E2
            iAE2[i] = iGenIrmsSJ[i]
            #########################################################################
            iStatorId[i] = -iR2[i] / 1.5 / G2
            iX2[i] = iStatorId[i]
            iStatorIq[i] = -D2 / 1.5 / G2
            iW2[i] = iStatorIq[i]
            iStatorIrms[i] = (iW2[i] ** 2 + iX2[i] ** 2) ** 0.5 / 1.414
            iY2[i] = iStatorIrms[i]
            iNetIrms[i] = 1 / M2 * iS2[i] / 1.732 / F2
            iAF2[i] = iNetIrms[i]
            #########################################################################
            iGenVd[i] = Rr * iGenId[i] - iZhuanChaLv[i] * (NetFreq * 2 * pi) * (
                    Lm * (iStatorIq[i]) + (Lkr + Lm) * iGenIq[i])
            iZ2[i] = iGenVd[i]
            iGenVq[i] = Rr * iGenIq[i] + iZhuanChaLv[i] * (NetFreq * 2 * pi) * (
                    Lm * (iStatorId[i]) + (Lm + Lkr) * iGenId[i])
            iAA2[i] = iGenVq[i]

            iGenV[i] = (iZ2[i] ** 2 + iAA2[i] ** 2) ** 0.5
            iAB2[i] = iGenV[i]
            iGenVrms[i] = iAB2[i] * E2 / F2 * 1.732 / 1.414
            iAC2 = iGenVrms[i]
            iTorque[i] = 1.5 * pp * Lm * (iGenId[i] * iStatorIq[i] - iGenIq[i] * iStatorId[i])

            iStatorAPower[i] = iStatorAPower[i] / 1E3
            iRotorAPower[i] = iRotorAPower[i] / -1E3

            iTotalAPower[i] = iStatorAPower[i] + iRotorAPower[i]
            iGridrms[i] = abs(iNetIrms[i]) + abs(iStatorIrms[i])

        # output part
        print('Torque = %f' % iTorque[-1])
        print('iU2 %f' % iU2[-1])
        print('iV2 %f' % iV2[-1])
        print('ZhuanChaLv iQ2 %f' % iQ2[-1])
        print('iW2 %f' % iW2[-1])
        print('iX2 %f' % iX2[-1])

        self.m_SynsSpeed.Value = ("%.3f" % TongBuZhuanSu)
        self.m_ZhuanChaLv.Value = ("%.3f" % iZhuanChaLv[-1])
        self.m_StatorAPower.Value = ("%.0f" % iStatorAPower[-1])  # 定子有功
        self.m_RotorAPower.Value = ("%.0f" % iRotorAPower[-1])  # 转子有功
        # self.m_RotorVoltRMS.Value = ("%8.3f" % iRotorVoltRMS[-1])  # 转子电压

        #########################################################################机网侧电流
        self.m_GenId.Value = ("%.3f" % iGenId[-1])  #
        self.m_GenIq.Value = ("%.3f" % iGenIq[-1])  #
        self.m_GenIrmsGS.Value = ("%.3f" % iGenIrmsGS[-1])  #
        self.m_GenIrmsSJ.Value = ("%.3f" % iGenIrmsSJ[-1])  #

        #########################################################################定子电流
        self.m_StatorId.Value = ("%.3f" % iStatorId[-1])  #
        self.m_StatorIq.Value = ("%.3f" % iStatorIq[-1])  #
        self.m_StatorIrms.Value = ("%.3f" % iStatorIrms[-1])  #
        self.m_NetIrms.Value = ("%.3f" % iNetIrms[-1])  #
        self.m_NetId.Value = ("%.3f" % (iNetIrms[-1] * 1.414))  #
        self.m_NetIq.Value = ("%.3f" % (0))  #

        #########################################################################电压扭矩
        self.m_GenVd.Value = ("%.3f" % iGenVd[-1])  #
        self.m_GenVq.Value = ("%.3f" % iGenVq[-1])  #
        self.m_GenV.Value = ("%.3f" % iGenV[-1])  #
        self.m_GenVrms.Value = ("%.3f" % iGenVrms[-1])  #
        self.m_Torque.Value = ("%.3f" % iTorque[-1])  #

        ######################################################################################

    def PlotAll(self):
        # pass
        plt.figure(5)
        gongLvTu = plt.subplot(2, 2, 1)  # 功率图
        gongLvTu.grid(linestyle='-.', which='major')
        gongLvTu.set_title('Scope - Power')

        miloc = plt.MultipleLocator(50)
        maloc = plt.MultipleLocator(100)
        gongLvTu.xaxis.set_minor_locator(miloc)
        gongLvTu.yaxis.set_minor_locator(maloc)
        gongLvTu.grid(linestyle='-.', which='minor', linewidth=0.3, alpha=0.9)

        iGenRpm_1 = iGenRpm[0:-1]
        iStatorAPower_1 = iStatorAPower[0:-1]
        iRotorAPower_1 = iRotorAPower[0:-1]
        iTotalAPower_1 = iTotalAPower[0:-1]

        plt.plot(iGenRpm_1, iStatorAPower_1, color="red", linewidth=1.5, linestyle="-", label="Stator Active Power")
        plt.plot(iGenRpm_1, iRotorAPower_1, color="blue", linewidth=1.5, linestyle="-.", label="Rotor Active Power")
        plt.plot(iGenRpm_1, iTotalAPower_1, color="gray", linewidth=1.5, linestyle="--", label="Total Active Power")

        plt.xlim(tempMin - 50, tempMax + 50)
        gongLvTu.set_xlabel('generetor speed[rpm]')
        gongLvTu.set_ylabel('[kW]')
        gongLvTu.legend()

        # 标注额定点
        for i in range(n - 1):
            gongLvTu.scatter(iGenRpm_1[i], (iStatorAPower_1[i]), s=30, color='gray')
            gongLvTu.scatter(iGenRpm_1[i], (iRotorAPower_1[i]), s=30, color='gray')
            gongLvTu.scatter(iGenRpm_1[i], (iTotalAPower_1[i]), s=30, color='gray')

        ######################################################################################
        dianLiuTu = plt.subplot(2, 2, 2)  # 电流图
        dianLiuTu.set_title('Scope - Current')
        dianLiuTu.grid(linestyle='-.', which='major')

        miloc = plt.MultipleLocator(50)
        maloc = plt.MultipleLocator(100)
        dianLiuTu.xaxis.set_minor_locator(miloc)
        dianLiuTu.yaxis.set_minor_locator(maloc)
        dianLiuTu.grid(linestyle='-.', which='minor', linewidth=0.3, alpha=0.9)

        iStatorIrms_1 = iStatorIrms[0:-1]
        iGenIrmsSJ_1 = iGenIrmsSJ[0:-1]
        iNetIrms_1 = iNetIrms[0:-1]
        iGridrms_1 = iGridrms[0:-1]

        plt.plot(iGenRpm_1, iStatorIrms_1, color="red", linewidth=1.5, linestyle="-", label="Istator RMS")  # 定子电流
        plt.plot(iGenRpm_1, iGenIrmsSJ_1, color="blue", linewidth=1.5, linestyle="--", label="Irotor RMS")  # 机侧电流
        plt.plot(iGenRpm_1, iNetIrms_1, color="gray", linewidth=1.5, linestyle="-.", label="Ipfc RMS")  # 网侧电流
        plt.plot(iGenRpm_1, iGridrms_1, color="green", linewidth=1.5, linestyle=":", label="Igrid_total RMS")  # 上网电流

        plt.xlim(tempMin - 50, tempMax + 50)
        dianLiuTu.set_xlabel('generetor speed[rpm]')
        dianLiuTu.set_ylabel('[A]')
        plt.legend()
        # 标注额定点
        for i in range(n - 1):
            dianLiuTu.scatter(iGenRpm_1[i], (iStatorIrms_1[i]), s=30, color='gray')
            dianLiuTu.scatter(iGenRpm_1[i], (iGenIrmsSJ_1[i]), s=30, color='gray')
            dianLiuTu.scatter(iGenRpm_1[i], (iNetIrms_1[i]), s=30, color='gray')
            dianLiuTu.scatter(iGenRpm_1[i], (iGridrms_1[i]), s=30, color='gray')

        ######################################################################################
        zhuanJuTu = plt.subplot(2, 2, 3)  # 转矩图
        zhuanJuTu.set_title('Scope - Torque')
        zhuanJuTu.grid(linestyle='-.', which='major')

        miloc = plt.MultipleLocator(50)
        maloc = plt.MultipleLocator(1000)
        zhuanJuTu.xaxis.set_minor_locator(miloc)
        zhuanJuTu.yaxis.set_minor_locator(maloc)
        zhuanJuTu.grid(linestyle='-.', which='minor', linewidth=0.3, alpha=0.9)

        iTorque_1 = iTorque[0:-1]

        plt.plot(iGenRpm_1, iTorque_1, color="red", linewidth=1.5, linestyle="-", label="Te")  # 电磁转矩

        plt.xlim(tempMin - 50, tempMax + 50)
        zhuanJuTu.set_xlabel('generetor speed[rpm]')
        zhuanJuTu.set_ylabel('[N.m]')
        plt.legend()
        # 标注额定点
        for i in range(n - 1):
            zhuanJuTu.scatter(iGenRpm_1[i], (iTorque_1[i]), s=30, color='gray')

        ######################################################################################
        dianYaTu = plt.subplot(2, 2, 4)  # 电压图
        dianYaTu.set_title('Scope - Rotor Voltage')
        dianYaTu.grid(linestyle='-.', which='major')

        miloc = plt.MultipleLocator(50)
        maloc = plt.MultipleLocator(100)
        dianYaTu.xaxis.set_minor_locator(miloc)
        dianYaTu.yaxis.set_minor_locator(maloc)
        dianYaTu.grid(linestyle='-.', which='minor', linewidth=0.3, alpha=0.9)

        iRotorVoltRMS_1 = iRotorVoltRMS[0:-1]
        iGenVrms_1 = iGenVrms[0:-1]

        # plt.plot(iGenRpm, iRotorVoltRMS, color="red", linewidth=1.5, linestyle="-", label="Vr RMS (estimate)")  # 定子电流
        plt.plot(iGenRpm_1, iGenVrms_1, color="red", linewidth=1.5, linestyle="-", label="Vinv RMS")  # 定子电流

        plt.xlim(tempMin - 50, tempMax + 50)
        dianYaTu.set_xlabel('generetor speed[rpm]')
        dianYaTu.set_ylabel('[V]')
        plt.legend()
        # 标注额定点
        for i in range(n - 1):
            # dianYaTu.scatter(iGenRpm[i], (iRotorVoltRMS[i]), s=30, color='gray')
            dianYaTu.scatter(iGenRpm_1[i], (iGenVrms_1[i]), s=30, color='gray')

        plt.tight_layout()
        # plt.savefig("fig1_4" + ".png", dpi=1000)
        plt.show()
        plt.close('all')

    def Plot1(self):
        # pass
        iGenRpm_1 = iGenRpm[0:-1]
        plt.figure(1)
        gongLvTu = plt.subplot(1, 1, 1)  # 功率图
        gongLvTu.grid(linestyle='-.', which='major')
        gongLvTu.set_title('Scope - Power')

        miloc = plt.MultipleLocator(50)
        maloc = plt.MultipleLocator(100)
        gongLvTu.xaxis.set_minor_locator(miloc)
        gongLvTu.yaxis.set_minor_locator(maloc)
        gongLvTu.grid(linestyle='-.', which='minor', linewidth=0.3, alpha=0.9)

        iStatorAPower_1 = iStatorAPower[0:-1]
        iRotorAPower_1 = iRotorAPower[0:-1]
        iTotalAPower_1 = iTotalAPower[0:-1]

        plt.plot(iGenRpm_1, iStatorAPower_1, color="red", linewidth=1.5, linestyle="-", label="Stator Active Power")
        plt.plot(iGenRpm_1, iRotorAPower_1, color="blue", linewidth=1.5, linestyle="-.", label="Rotor Active Power")
        plt.plot(iGenRpm_1, iTotalAPower_1, color="gray", linewidth=1.5, linestyle="--", label="Total Active Power")

        plt.xlim(tempMin - 50, tempMax + 50)
        gongLvTu.set_xlabel('generetor speed[rpm]')
        gongLvTu.set_ylabel('[kW]')
        gongLvTu.legend()

        # 标注额定点
        for i in range(n - 1):
            gongLvTu.scatter(iGenRpm_1[i], (iStatorAPower_1[i]), s=30, color='gray')
            gongLvTu.scatter(iGenRpm_1[i], (iRotorAPower_1[i]), s=30, color='gray')
            gongLvTu.scatter(iGenRpm_1[i], (iTotalAPower_1[i]), s=30, color='gray')

        plt.tight_layout()

        plt.savefig(outputFolderName + r"\fig_power.png", dpi=1000)
        plt.close(1)

    def Plot2(self):
        # pass
        iGenRpm_1 = iGenRpm[0:-1]
        plt.figure(2)
        dianYaTu = plt.subplot(1, 1, 1)  # 电压图
        dianYaTu.set_title('Scope - Rotor Voltage')
        dianYaTu.grid(linestyle='-.', which='major')

        miloc = plt.MultipleLocator(50)
        maloc = plt.MultipleLocator(100)
        dianYaTu.xaxis.set_minor_locator(miloc)
        dianYaTu.yaxis.set_minor_locator(maloc)
        dianYaTu.grid(linestyle='-.', which='minor', linewidth=0.3, alpha=0.9)

        iRotorVoltRMS_1 = iRotorVoltRMS[0:-1]
        iGenVrms_1 = iGenVrms[0:-1]

        # plt.plot(iGenRpm, iRotorVoltRMS, color="red", linewidth=1.5, linestyle="-", label="Vr RMS (estimate)")  # 定子电流
        plt.plot(iGenRpm_1, iGenVrms_1, color="red", linewidth=1.5, linestyle="-", label="Vinv RMS")  # 定子电流

        plt.xlim(tempMin - 50, tempMax + 50)
        dianYaTu.set_xlabel('generetor speed[rpm]')
        dianYaTu.set_ylabel('[V]')
        plt.legend()
        # 标注额定点
        for i in range(n - 1):
            # dianYaTu.scatter(iGenRpm[i], (iRotorVoltRMS[i]), s=30, color='gray')
            dianYaTu.scatter(iGenRpm_1[i], (iGenVrms_1[i]), s=30, color='gray')

        plt.tight_layout()
        plt.savefig(outputFolderName + r"\fig_voltageR.png", dpi=1000)
        plt.close(2)

    def Plot3(self):
        iGenRpm_1 = iGenRpm[0:-1]
        plt.figure(3)
        dianLiuTu = plt.subplot(1, 1, 1)  # 电流
        dianLiuTu.set_title('Scope - Current')
        dianLiuTu.grid(linestyle='-.', which='major')

        miloc = plt.MultipleLocator(50)
        maloc = plt.MultipleLocator(100)
        dianLiuTu.xaxis.set_minor_locator(miloc)
        dianLiuTu.yaxis.set_minor_locator(maloc)
        dianLiuTu.grid(linestyle='-.', which='minor', linewidth=0.3, alpha=0.9)

        iStatorIrms_1 = iStatorIrms[0:-1]
        iGenIrmsSJ_1 = iGenIrmsSJ[0:-1]
        iNetIrms_1 = iNetIrms[0:-1]
        iGridrms_1 = iGridrms[0:-1]

        plt.plot(iGenRpm_1, iStatorIrms_1, color="red", linewidth=1.5, linestyle="-", label="Istator RMS")  # 定子电流
        plt.plot(iGenRpm_1, iGenIrmsSJ_1, color="blue", linewidth=1.5, linestyle="--", label="Irotor RMS")  # 机侧电流
        plt.plot(iGenRpm_1, iNetIrms_1, color="gray", linewidth=1.5, linestyle="-.", label="Ipfc RMS")  # 网侧电流
        plt.plot(iGenRpm_1, iGridrms_1, color="green", linewidth=1.5, linestyle=":", label="Igrid_total RMS")  # 上网电流

        plt.xlim(tempMin - 50, tempMax + 50)
        dianLiuTu.set_xlabel('generetor speed[rpm]')
        dianLiuTu.set_ylabel('[A]')
        plt.legend()
        # 标注额定点
        for i in range(n - 1):
            dianLiuTu.scatter(iGenRpm_1[i], (iStatorIrms_1[i]), s=30, color='gray')
            dianLiuTu.scatter(iGenRpm_1[i], (iGenIrmsSJ_1[i]), s=30, color='gray')
            dianLiuTu.scatter(iGenRpm_1[i], (iNetIrms_1[i]), s=30, color='gray')
            dianLiuTu.scatter(iGenRpm_1[i], (iGridrms_1[i]), s=30, color='gray')

        plt.tight_layout()
        plt.savefig(outputFolderName + r"\fig_Current.png", dpi=1000)
        plt.close(3)

    def Docu(self):

        '''save %0 data'''
        iTotalAPower100 = iTotalAPower
        iStatorAPower100 = iStatorAPower
        iRotorAPower100 = iRotorAPower
        iStatorIrms100 = list(map(abs, iStatorIrms))
        iGenIrmsSJ100 = list(map(abs, iGenIrmsSJ))
        iNetIrms100 = list(map(abs, iNetIrms))
        iGridrms100 = list(map(abs, iGridrms))
        iGenVrms100 = iGenVrms

        self.DoWork(1.1)
        iTotalAPower110 = iTotalAPower
        iStatorAPower110 = iStatorAPower
        iRotorAPower110 = iRotorAPower
        iStatorIrms110 = list(map(abs, iStatorIrms))
        iGenIrmsSJ110 = list(map(abs, iGenIrmsSJ))
        iNetIrms110 = list(map(abs, iNetIrms))
        iGridrms110 = list(map(abs, iGridrms))
        iGenVrms110 = iGenVrms

        self.DoWork(0.9)
        iTotalAPower90 = iTotalAPower
        iStatorAPower90 = iStatorAPower
        iRotorAPower90 = iRotorAPower
        iStatorIrms90 = list(map(abs, iStatorIrms))
        iGenIrmsSJ90 = list(map(abs, iGenIrmsSJ))
        iNetIrms90 = list(map(abs, iNetIrms))
        iGridrms90 = list(map(abs, iGridrms))
        iGenVrms90 = iGenVrms

        self.DoWork(1)
        #####################################################################创建PPT,修改文本框
        print('ppt start')
        pptx = Presentation()
        for slide in pptx.slides:
            # 遍历幻灯片页的所有形状
            print("slide found")
            for shape in slide.shapes:
                print("shape found")
                # 判断形状是否含有文本框，如果含有则顺序运行代码
                if shape.has_text_frame:

                    # 获取文本框
                    text_frame = shape.text_frame
                    # 遍历文本框中的所有段落
                    for paragraph in text_frame.paragraphs:
                        paragraph.text = paragraph.text.replace('X1', str(format(Xs, '.5f')) + '[ohm]')
                        paragraph.text = paragraph.text.replace('X2', str(format(Xr, '.5f')) + '[ohm]')
                        paragraph.text = paragraph.text.replace('R1', str(format(Rs, '.5f')) + '[ohm]')
                        paragraph.text = paragraph.text.replace('R2', str(format(Rr, '.5f')) + '[ohm]')

                        paragraph.text = paragraph.text.replace('XNN', str(format(Xm, '.5f')) + '[ohm]')
                        paragraph.font.size = Pt(22)

        #####################################################################保存PPT,另存为png
        save_ppt = outputFolderName + r'\fig_x.pptx'
        pptx.save(save_ppt)

        utils.save_pptx_as_png(outputFolderName, save_ppt, overwrite_folder=True)
        fig_x_dir = outputFolderName + r'\幻灯片1.PNG'
        print('ppt end')
        #####################################################################创建word
        document = Document()
        # print(dir(document))

        document.styles['Normal'].font.name = '楷体'

        document.styles['Normal']._element.rPr.rFonts.set(qn('w:eastAsia'), u'楷体')
        # document.styles['Normal'].font.size = Pt(12)
        document.add_heading('双馈风电机组电驱系统设计计算报告', 0)

        p = document.add_paragraph('生成时间   ' + datetime.datetime.now().strftime('%Y.%m.%d  %H:%M'))
        # p.add_run('bold').bold = True
        # p.add_run(' and some ')
        # p.add_run('italic.').italic = True

        document.add_heading('输入-电机参数', level=1)
        # document.add_paragraph(
        #      '极对数', style='List Number'
        #  )
        ###################参数表######################
        table_para = document.add_table(rows=12, cols=2)
        # print(dir(table_para))
        hdr_cells = table_para.columns[0].cells

        hdr_cells[0].text = '极对数'
        hdr_cells[1].text = '上网总功率[kW]'
        hdr_cells[2].text = '转子开口电压[Vrms]'
        hdr_cells[3].text = '电网电压[Vrms]'
        hdr_cells[4].text = '电网频率[Hz]'
        hdr_cells[5].text = '定子无功功率[kW]'
        hdr_cells[6].text = '变流器效率'
        hdr_cells[7].text = '直流母线电压[V]'
        hdr_cells[8].text = '网侧变流器最大电流[A]'
        hdr_cells[9].text = '机侧变流器最大电流[A]'
        hdr_cells[10].text = '发电机定子最大电流[A]'
        hdr_cells[11].text = ' '

        para_cells = table_para.columns[1].cells
        para_cells[0].text = self.m_pp.GetValue()
        para_cells[1].text = self.m_GoNetPower.GetValue()
        para_cells[2].text = self.m_RotorOpen.GetValue()
        para_cells[3].text = self.m_GridVRMS.GetValue()
        para_cells[4].text = self.m_hz.GetValue()
        para_cells[5].text = self.m_StatorRePower.GetValue()
        para_cells[6].text = self.m_XiaoLv.GetValue()
        para_cells[7].text = self.m_Vdc.GetValue()
        para_cells[8].text = self.m_SetNetCurMax.GetValue()
        para_cells[9].text = self.m_SetGenCurMax1.GetValue()
        para_cells[10].text = self.m_SetSatMaxI.GetValue()
        para_cells[11].text = ' '

        ###################电机表######################
        table_para = document.add_table(rows=5, cols=2)
        # print(dir(table_para))
        hdr_cells = table_para.columns[0].cells

        self.m_Lksss.SetLabel(u"定子漏抗[ohm]")
        self.m_Lkrrr.SetLabel(u"转子漏抗[ohm]")
        self.m_staticText22.SetLabel(u"激磁电抗[ohm]")

        hdr_cells[0].text = self.m_staticText22.GetLabel()
        hdr_cells[1].text = '定子电组[ohm]'
        hdr_cells[2].text = '转子电阻[ohm]'
        hdr_cells[3].text = self.m_Lksss.GetLabel()
        hdr_cells[4].text = self.m_Lkrrr.GetLabel()

        para_cells = table_para.columns[1].cells
        para_cells[0].text = self.m_Lm.GetValue()
        para_cells[1].text = self.m_Rs.GetValue()
        para_cells[2].text = self.m_Rr.GetValue()
        para_cells[3].text = self.m_Lks.GetValue()
        para_cells[4].text = self.m_Lkr.GetValue()

        ########等效电路图
        try:
            paragraph = document.add_paragraph()
            paragraph.alignment = WD_PARAGRAPH_ALIGNMENT.CENTER
            run = paragraph.add_run("")
            run.add_picture(r'.\fig_save\幻灯片1.PNG', width=Inches(5.25))

            run = document.add_paragraph('图1. 双馈发电机等效电路图')
            run.alignment = WD_PARAGRAPH_ALIGNMENT.CENTER

        except:
            pass

        ########功率表
        document.add_heading('输入-功率曲线', level=1)
        n = len(iGenRpm) - 1
        table = document.add_table(rows=1, cols=3)
        hdr_cells = table.rows[0].cells
        hdr_cells[0].text = '转速[rpm]'
        hdr_cells[1].text = '上网功率[kW]'
        hdr_cells[2].text = ' '

        for i in range(n):
            row_cells = table.add_row().cells
            row_cells[0].text = str(iGenRpm[i])

            row_cells[1].text = str(format(iGoNetPower[i] * 0.001, '.2f'))
            row_cells[2].text = ' '

        ########输出功率
        document.add_heading('输出-功率计算', level=1)

        paragraph = document.add_paragraph()
        paragraph.alignment = WD_PARAGRAPH_ALIGNMENT.CENTER
        run = paragraph.add_run("")
        run.add_picture(r'.\fig_save\fig_power.png', width=Inches(5.25))

        run = document.add_paragraph('图2. 定子/转子/总上网功率')
        run.alignment = WD_PARAGRAPH_ALIGNMENT.CENTER
        run = document.add_paragraph('')
        run.add_run(r'"*注：对于转子功率(Rotor Active Power)，负值表示从电网获取功率"').italic = True
        document.add_paragraph('总上网功率最大值为：\t%8.3f [kW]' % (max(iTotalAPower100)))
        document.add_paragraph('定子功率最大值为：  \t%8.3f [kW]' % (max(iStatorAPower100)))
        document.add_paragraph('转子功率最大值为：  \t%8.3f [kW]' % (max(iRotorAPower100)))
        document.add_paragraph('转子功率最小值为：  \t%8.3f [kW]' % (min(iRotorAPower100)))
        document.add_page_break()

        ########输出电流
        document.add_heading('输出-电流', level=1)

        paragraph = document.add_paragraph()
        paragraph.alignment = WD_PARAGRAPH_ALIGNMENT.CENTER
        run = paragraph.add_run("")
        run.add_picture(r'.\fig_save\fig_Current.png', width=Inches(5.25))

        run = document.add_paragraph('图3. 定子/转子/变流器网侧/总上网电流')
        run.alignment = WD_PARAGRAPH_ALIGNMENT.CENTER
        run = document.add_paragraph('')
        run.add_run(r'"*注：对于变流器网侧电流(Ipfc)，负值表示流向电网"').italic = True
        run.add_run('\n').italic = True
        run.add_run(r'"*注：转子电流、变流器网侧电流均为实际值，而非归算值"').italic = True

        document.add_paragraph(r'以下三个数据分别表示工况（正常电网电压）（电网90%低电压）（电网110%高电压）')
        document.add_paragraph('定子电流最大值为：      \t%8.3f [A],   %8.3f [A],   %8.3f [A]' % (
            max(iStatorIrms100), max(iStatorIrms90), max(iStatorIrms110)))
        document.add_paragraph('转子电流最大值为：      \t%8.3f [A],   %8.3f [A],   %8.3f [A]' % (
            max(iGenIrmsSJ100), max(iGenIrmsSJ90), max(iGenIrmsSJ110)))
        document.add_paragraph(
            '变流器网侧电流最大值为： \t%8.3f [A],   %8.3f [A],   %8.3f [A]' % (
                max(iNetIrms100), max(iNetIrms90), max(iNetIrms110)))
        document.add_paragraph(
            '总上网电流最大值为：    \t%8.3f [A],   %8.3f [A],   %8.3f [A]' % (
                max(iGridrms100), max(iGridrms90), max(iGridrms110)))
        document.add_page_break()

        ########输出转子电压
        document.add_heading('输出-转子电压', level=1)

        paragraph = document.add_paragraph()
        paragraph.alignment = WD_PARAGRAPH_ALIGNMENT.CENTER
        run = paragraph.add_run("")
        run.add_picture(r'.\fig_save\fig_voltageR.png', width=Inches(5.25))

        run = document.add_paragraph('图4. 转子线电压')
        run.alignment = WD_PARAGRAPH_ALIGNMENT.CENTER

        document.add_paragraph(r'以下三个数据分别表示工况（正常电网电压）（电网90%低电压）（电网110%高电压）')
        document.add_paragraph(
            '转子线电压最大值为： \t%8.3f [V],   %8.3f [V],   %8.3f [V]' % (max(iGenVrms100), max(iGenVrms90), max(iGenVrms110)))
        document.add_page_break()

        ########结论
        document.add_heading('结论', level=1)
        table_para = document.add_table(rows=5, cols=3, style='Table Grid')

        hdr_cells = table_para.columns[0].cells

        hdr_cells[0].text = '判定项'
        hdr_cells[1].text = '发电机转子电压'
        hdr_cells[2].text = '网侧变流器电流'
        hdr_cells[3].text = '机侧变流器电流'
        hdr_cells[4].text = '发电机定子电流'

        para_cells1 = table_para.columns[1].cells
        para_cells1[0].text = '判定标准'

        TF = [True for i in range(4)]
        para_cells1[1].text = str(format(max(max(iGenVrms90), max(iGenVrms100), max(iGenVrms110)),
                                         '.2f')) + ' < ' + self.m_Vdc.GetValue() + '/1.414'
        TF[0] = max(max(iGenVrms90), max(iGenVrms100), max(iGenVrms110)) < float(self.m_Vdc.GetValue()) / 1.414

        para_cells1[2].text = str(format(max(max(iNetIrms90), max(iNetIrms100), max(iNetIrms110)),
                                         '.2f')) + ' < ' + self.m_SetNetCurMax.GetValue()
        TF[1] = max(max(iNetIrms90), max(iNetIrms100), max(iNetIrms110)) < float(self.m_SetNetCurMax.GetValue())

        para_cells1[3].text = str(format(max(max(iGenIrmsSJ90), max(iGenIrmsSJ100), max(iGenIrmsSJ110)),
                                         '.2f')) + ' < ' + self.m_SetGenCurMax1.GetValue()
        TF[2] = max(max(iGenIrmsSJ90), max(iGenIrmsSJ100), max(iGenIrmsSJ110)) < float(self.m_SetGenCurMax1.GetValue())

        para_cells1[4].text = str(format(max(max(iStatorIrms90), max(iStatorIrms100), max(iStatorIrms110)),
                                         '.2f')) + ' < ' + self.m_SetSatMaxI.GetValue()
        TF[3] = max(max(iStatorIrms90), max(iStatorIrms100), max(iStatorIrms110)) < float(self.m_SetSatMaxI.GetValue())

        para_cells2 = table_para.columns[2].cells
        para_cells2[0].text = '是否通过'

        #1,2,3,4行
        for i in range(4):
            p = para_cells2[i + 1].paragraphs.pop()
            #print(dir(para_cells2[i + 1].paragraphs))
            #print(dir(p.style))
            #print(len(p))
            #p.style.font.color.rgb = RGBColor(250, 0, 0)
            run = p.add_run(str(TF[i]))
            if TF[i] is False:
                run.font.color.rgb = RGBColor(250, 0, 0)



        document.add_paragraph(r'')

        if all(i == True for i in TF):
            document.add_paragraph(r'结论:设计参数校验通过。')
        else:
            paragraph = document.add_paragraph(r'结论:设计参数 %d项校验不通过！' % TF.count(False))
            # paragraph.style.font.color.rgb = RGBColor(250,0,0)

        document.add_page_break()
        document.save('双馈风电机组电驱系统设计计算报告.docx')
        print(TF.count(False))
        print(TF)

        ##############################################test

    def DocuTemp(self):
        tdocx = Document('tpl.docx')
        tdocx.styles['Normal'].font.name = '微软雅黑'
        tdocx.styles['Normal']._element.rPr.rFonts.set(qn('w:eastAsia'), u'微软雅黑')
        # tdocx.styles['Normal'].font.size = Pt(12)

        tdocx.add_page_break()
        tdocx.save('test.docx')

        ##############################################test


########################################################## 3.45 盾安 MW

# self.m_grid2.SetCellValue(0, 0, '1300')  # 第1行第0个
# self.m_grid2.SetCellValue(0, 1, '1400')
# self.m_grid2.SetCellValue(0, 2, '1450')
# self.m_grid2.SetCellValue(0, 3, '1550')
# self.m_grid2.SetCellValue(0, 4, '1600')
# self.m_grid2.SetCellValue(0, 5, '1650')
# self.m_grid2.SetCellValue(0, 6, '1700')
# self.m_grid2.SetCellValue(0, 7, '1750')
# self.m_grid2.SetCellValue(0, 8, '1800')
#
# self.m_grid2.SetCellValue(1, 0, '1454')  # 第1行第0个
# self.m_grid2.SetCellValue(1, 1, '1681')
# self.m_grid2.SetCellValue(1, 2, '1851')
# self.m_grid2.SetCellValue(1, 3, '2151')
# self.m_grid2.SetCellValue(1, 4, '2352')
# self.m_grid2.SetCellValue(1, 5, '2502')
# self.m_grid2.SetCellValue(1, 6, '2710')
# self.m_grid2.SetCellValue(1, 7, '3153')
# self.m_grid2.SetCellValue(1, 8, '3452')

########################################################## 3.35 MW

# self.m_grid2.SetCellValue(0, 0, '1050')  # 第1行第0个
# self.m_grid2.SetCellValue(0, 1, '1200')
# self.m_grid2.SetCellValue(0, 2, '1300')
# self.m_grid2.SetCellValue(0, 3, '1400')
# self.m_grid2.SetCellValue(0, 4, '1550')
# self.m_grid2.SetCellValue(0, 5, '1600')
# self.m_grid2.SetCellValue(0, 6, '1700')
# self.m_grid2.SetCellValue(0, 7, '1725')
# self.m_grid2.SetCellValue(0, 8, '1770')
#
# self.m_grid2.SetCellValue(1, 0, '400')  # 第1行第0个
# self.m_grid2.SetCellValue(1, 1, '700')
# self.m_grid2.SetCellValue(1, 2, '850')
# self.m_grid2.SetCellValue(1, 3, '1150')
# self.m_grid2.SetCellValue(1, 4, '1500')
# self.m_grid2.SetCellValue(1, 5, '1750')
# self.m_grid2.SetCellValue(1, 6, '2250')
# self.m_grid2.SetCellValue(1, 7, '3250')
# self.m_grid2.SetCellValue(1, 8, '3350')

########################################################## 3.0 鼠笼MW
# self.m_grid2.SetCellValue(0, 0, '660.568')  # 第1行第0个
# self.m_grid2.SetCellValue(0, 1, '779.144')
# self.m_grid2.SetCellValue(0, 2, '896.907')
# self.m_grid2.SetCellValue(0, 3, '1013.69')
# self.m_grid2.SetCellValue(0, 4, '1129.12')
# self.m_grid2.SetCellValue(0, 5, '1243.81')
# self.m_grid2.SetCellValue(0, 6, '1357.99')
# self.m_grid2.SetCellValue(0, 7, '1470.47')
# self.m_grid2.SetCellValue(0, 8, '1583.24')
# self.m_grid2.SetCellValue(0, 9, '1693.32')
# self.m_grid2.SetCellValue(0, 10, '1725')
# self.m_grid2.SetCellValue(0, 11, '1725')
#
# self.m_grid2.SetCellValue(1, 0, '59.2634')  # 第1行第0个
# self.m_grid2.SetCellValue(1, 1, '126.16')
# self.m_grid2.SetCellValue(1, 2, '216.215')
# self.m_grid2.SetCellValue(1, 3, '333.262')
# self.m_grid2.SetCellValue(1, 4, '478.754')
# self.m_grid2.SetCellValue(1, 5, '648.541')
# self.m_grid2.SetCellValue(1, 6, '851.05')
# self.m_grid2.SetCellValue(1, 7, '1087.73')
# self.m_grid2.SetCellValue(1, 8, '1366.06')
# self.m_grid2.SetCellValue(1, 9, '1676.12')
# self.m_grid2.SetCellValue(1, 10, '1819')
# self.m_grid2.SetCellValue(1, 11, '3000')
# self.m_grid2.SetCellValue(1, 11, '3000')
'''end of file'''
